"""Generate binary fixture files for parser tests. Run once from backend/."""
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def _make_pdf(text: str) -> bytes:
    """Build a minimal single-page PDF with extractable text (no Rust/reportlab needed)."""
    content = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()

    objs: list[bytes] = []
    objs.append(b"1 0 obj\n<</Type /Catalog /Pages 2 0 R>>\nendobj\n")
    objs.append(b"2 0 obj\n<</Type /Pages /Kids [3 0 R] /Count 1>>\nendobj\n")
    objs.append(
        b"3 0 obj\n"
        b"<</Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]\n"
        b"  /Resources <</Font <</F1 5 0 R>>>> /Contents 4 0 R>>\n"
        b"endobj\n"
    )
    stream_body = b"stream\n" + content + b"\nendstream\n"
    objs.append(
        f"4 0 obj\n<</Length {len(content)}>>\n".encode() + stream_body + b"endobj\n"
    )
    objs.append(
        b"5 0 obj\n"
        b"<</Type /Font /Subtype /Type1 /BaseFont /Helvetica>>\n"
        b"endobj\n"
    )

    header = b"%PDF-1.4\n"
    offsets: list[int] = []
    body = b""
    pos = len(header)
    for obj in objs:
        offsets.append(pos)
        body += obj
        pos += len(obj)

    xref_start = pos
    xref = b"xref\n"
    xref += f"0 {len(objs) + 1}\n".encode()
    xref += b"0000000000 65535 f \n"
    for o in offsets:
        xref += f"{o:010d} 00000 n \n".encode()

    trailer = (
        f"trailer\n<</Size {len(objs) + 1} /Root 1 0 R>>\n"
        f"startxref\n{xref_start}\n%%EOF\n"
    ).encode()

    return header + body + xref + trailer


def main() -> None:
    # PDF
    (FIXTURES_DIR / "sample.pdf").write_bytes(
        _make_pdf(
            "This is a sample PDF document for testing the parser. "
            "It contains a full sentence with enough characters."
        )
    )

    # DOCX
    from docx import Document as DocxDoc

    doc = DocxDoc()
    doc.add_heading("Introduction", level=1)
    doc.add_paragraph(
        "This is a sample DOCX document for testing the parser. "
        "It contains enough text to produce at least one chunk."
    )
    doc.add_heading("Main Content", level=1)
    doc.add_paragraph(
        "Here is the main content section. "
        "The parser should identify this heading and the text below it."
    )
    doc.save(FIXTURES_DIR / "sample.docx")

    # PPTX
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Sample Presentation"
    slide.placeholders[1].text = (
        "This is slide 1 content for testing the PPTX parser. "
        "It has enough text to produce at least one chunk."
    )
    prs.save(FIXTURES_DIR / "sample.pptx")

    # TXT
    (FIXTURES_DIR / "sample.txt").write_text(
        "This is a sample plain-text document.\n\n"
        "It has multiple paragraphs so the splitter has something to work with.\n\n"
        "A third paragraph ensures there is plenty of content for the tests.",
        encoding="utf-8",
    )

    # MD
    (FIXTURES_DIR / "sample.md").write_text(
        "# Sample Markdown\n\n"
        "This is a sample Markdown document for testing the text parser.\n\n"
        "## Section Two\n\n"
        "More content here ensures the parser works on markdown files too.",
        encoding="utf-8",
    )

    print(f"Fixtures written to {FIXTURES_DIR.resolve()}")


if __name__ == "__main__":
    main()
