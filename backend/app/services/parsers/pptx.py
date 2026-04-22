import logging

from pptx import Presentation

from app.services.parsers import ParsedChunk

logger = logging.getLogger(__name__)


def parse(file_path: str) -> list[ParsedChunk]:
    """Extract text slide by slide. page = slide number (1-indexed).

    Text is collected from all shapes and from speaker notes.
    Slides with no extractable text are skipped.
    """
    prs = Presentation(file_path)
    chunks: list[ParsedChunk] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")

        combined = "\n".join(parts)
        if combined.strip():
            chunks.append(ParsedChunk(text=combined, page=slide_num, section=None))
        else:
            logger.debug("Slide %d has no extractable text, skipping", slide_num)

    return chunks
